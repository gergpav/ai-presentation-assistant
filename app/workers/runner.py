# app/workers/runner.py
# КРИТИЧНО: Устанавливаем переменные окружения ДО импорта PyTorch/transformers
import os

# Принудительно отключаем CUDA (по умолчанию true для избежания ошибок)
force_cpu = os.getenv("FORCE_CPU", "true").lower() == "true"
if force_cpu:
    os.environ["CUDA_VISIBLE_DEVICES"] = ""  # Скрываем GPU от всех библиотек
    os.environ["TOKENIZERS_PARALLELISM"] = "false"

import asyncio
import logging
import uuid
from sqlalchemy.orm import Session, joinedload
from pathlib import Path
from typing import Optional, List

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
import pandas as pd

from sqlalchemy import select, desc
from sqlalchemy.ext.asyncio import AsyncSession

from app.db.session import AsyncSessionLocal
from app.db.models.job import Job
from app.db.models.project import Project
from app.db.models.slide import Slide
from app.db.models.slide_content import SlideContent
from app.db.models.slide_document import SlideDocument
from app.db.models.file import File
from app.db.models.enums import JobStatus, JobType, SlideStatus, FileKind
from app.db.models import Template

from app.core.llm_generator import content_generator
from app.core.embeddings import DocumentIndex

from app.core.pptx_builder import PresentationBuilder
from app.core.pdf_builder import slides_to_pdf_bytes
from app.core.image_generator import image_generator
from app.core.visual_generator import generate_table_image, generate_chart_image
from app.utils.helpers import SlideExport
from app.config import settings

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)

POLL_INTERVAL_SEC = 1.0
STORAGE_DIR = Path("storage")
STORAGE_DIR.mkdir(exist_ok=True)
IMAGES_DIR = STORAGE_DIR / "images"
IMAGES_DIR.mkdir(exist_ok=True)

# Глобальный семафор для ограничения параллельных вызовов генерации модели
# Инициализируется в worker_loop()
_llm_generation_semaphore: Optional[asyncio.Semaphore] = None


# ----------------------------
# Helpers: SlideContent text field (под разные имена колонок)
# ----------------------------
def _sc_get_text(sc: SlideContent | None) -> str:
    if sc is None:
        return ""
    if hasattr(sc, "content_text"):
        return sc.content_text or ""
    if hasattr(sc, "content"):
        # SlideContent.content это JSON поле (dict)
        if isinstance(sc.content, dict):
            return sc.content.get("text", "") or ""
        elif isinstance(sc.content, str):
            return sc.content or ""
        return ""
    if hasattr(sc, "text"):
        return sc.text or ""
    return ""


def _sc_set_text(sc: SlideContent, text: str) -> None:
    """Сохраняет текст в SlideContent.content как JSON"""
    if hasattr(sc, "content_text"):
        sc.content_text = text
        return
    if hasattr(sc, "content"):
        # SlideContent.content это JSON поле, сохраняем как словарь
        sc.content = {"text": text}
        return
    if hasattr(sc, "text"):
        sc.text = text
        return
    raise RuntimeError("SlideContent: не найдено поле для текста (content_text/content/text)")


# ----------------------------
# Helpers: parsing files from path (pptx/docx/xlsx/pdf)
# ----------------------------
async def _parse_text_from_path(path: str) -> str:
    """
    Мини-парсер "как есть" для контекста.
    Хранит итоговый текст в SlideDocument.parsed_text, чтобы не парсить повторно.
    """
    p = Path(path)
    if not p.exists():
        return ""

    ext = p.suffix.lower()

    try:
        if ext == ".pdf":
            text = ""
            with pdfplumber.open(str(p)) as pdf:
                for page in pdf.pages:
                    text += (page.extract_text() or "") + "\n"
            return text.strip()

        if ext == ".docx":
            doc = DocxDocument(str(p))
            parts = [par.text for par in doc.paragraphs if par.text]
            return "\n".join(parts).strip()

        if ext == ".pptx":
            pres = PptxPresentation(str(p))
            parts: list[str] = []
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            return "\n".join(parts).strip()

        if ext in (".xlsx", ".xls"):
            # читаем все листы, склеиваем как текстовую таблицу
            xls = pd.ExcelFile(str(p))
            out: list[str] = []
            for sheet in xls.sheet_names:
                df = xls.parse(sheet)
                out.append(f"=== {sheet} ===")
                out.append(df.to_string(index=False))
            return "\n".join(out).strip()

        # fallback: просто читаем как текст (если вдруг txt)
        return p.read_text(encoding="utf-8", errors="ignore").strip()

    except Exception as e:
        logger.warning(f"Failed to parse {path}: {e}")
        return ""


# ----------------------------
# Helpers: build context for a slide from SlideDocument
# ----------------------------
async def _build_slide_context(db: AsyncSession, slide_id: int, query: str) -> str:
    res = await db.execute(
        select(SlideDocument)
        .where(SlideDocument.slide_id == slide_id)
        .order_by(SlideDocument.id.asc())
    )
    docs = res.scalars().all()
    if not docs:
        return ""

    documents_for_index: list[dict] = []

    for d in docs:
        if d.parsed_text:
            text = d.parsed_text
        else:
            text = await _parse_text_from_path(d.storage_path)
            if text:
                d.parsed_text = text  # кешируем
        documents_for_index.append({"text": text, "metadata": {"filename": d.filename}})

    await db.commit()

    idx = DocumentIndex()
    idx.add_documents(documents_for_index)
    idx.build_index()

    hits = idx.search(query, k=5)
    context = "\n\n".join(f"[{source}]\n{content}" for (content, _kind, source) in hits if content)
    return context[:12000]


async def _get_latest_slide_content(db: AsyncSession, slide_id: int) -> SlideContent | None:
    res = await db.execute(
        select(SlideContent)
        .where(SlideContent.slide_id == slide_id)
        .order_by(desc(SlideContent.id))
        .limit(1)
    )
    return res.scalar_one_or_none()


# ----------------------------
# Jobs: export template (optional)
# ----------------------------
async def _resolve_project_template_path(db: AsyncSession, project_id: int) -> str | None:
    """
    Возвращает путь к pptx-шаблону проекта (если назначен).
    Поддерживает оба варианта:
      1) relationship: Project.template
      2) FK: Project.template_id
    """
    # 1) пробуем загрузить проект сразу с template (если relationship есть)
    try:
        res = await db.execute(
            select(Project)
            .options(joinedload(Project.template))
            .where(Project.id == project_id)
        )
        project = res.scalar_one()
        tmpl = getattr(project, "template", None)
        if tmpl and getattr(tmpl, "storage_path", None):
            return tmpl.storage_path
    except Exception:
        # relationship может отсутствовать — ок, падаем на fallback
        pass

    # 2) fallback: template_id -> Template
    res = await db.execute(select(Project).where(Project.id == project_id))
    project = res.scalar_one()

    template_id = getattr(project, "template_id", None)
    if not template_id:
        return None

    res = await db.execute(select(Template).where(Template.id == template_id))
    tmpl = res.scalar_one_or_none()
    if not tmpl:
        return None

    return getattr(tmpl, "storage_path", None)



# ----------------------------
# Job handlers
# ----------------------------
async def generate_slide_job(db: AsyncSession, job: Job) -> None:
    if job.slide_id is None:
        raise ValueError("Job.slide_id is required for generate_slide")

    slide = (await db.execute(select(Slide).where(Slide.id == job.slide_id))).scalar_one()
    project = (await db.execute(select(Project).where(Project.id == slide.project_id))).scalar_one()

    # Логируем тип визуализации для отладки
    logger.info(f"Генерация слайда {slide.id}: visual_type={slide.visual_type.value}, prompt={slide.prompt[:100]}...")

    if not slide.prompt or not slide.prompt.strip():
        raise ValueError("Slide prompt is empty")

    job.progress = 10
    await db.commit()

    context_text = await _build_slide_context(db, slide.id, slide.prompt)

    job.progress = 30
    await db.commit()

    # Инициализируем переменные для изображения заранее
    generated_image_path = None
    
    # Таймаут генерации контента из конфигурации (определяем заранее для использования в метаданных)
    # Можно переопределить через переменную окружения LLM_GENERATION_TIMEOUT_SEC
    timeout_sec = settings.LLM_GENERATION_TIMEOUT_SEC
    
    # Проверяем, является ли это первым слайдом (титульный слайд)
    # или название слайда указывает на титульный слайд
    is_first_slide = slide.position == 1
    title_lower = (slide.title or "").lower()
    title_keywords = ["титульный", "титул", "обложка", "cover", "title slide", "начало"]
    is_title_by_name = any(kw in title_lower for kw in title_keywords)
    
    # Если это первый слайд или название указывает на титульный, принудительно устанавливаем layout = "title"
    if is_first_slide and is_title_by_name:
        layout_type = "title"
        # Для титульного слайда генерируем заголовок и подзаголовок из промпта
        # (не используем название слайда, так как оно может быть просто "Титульный слайд")
        job.progress = 40
        await db.commit()
        
        try:
            logger.info(f"Генерация титульного слайда (таймаут: {timeout_sec}с)")
            if _llm_generation_semaphore:
                async with _llm_generation_semaphore:
                    out = await asyncio.wait_for(
                        asyncio.to_thread(
                            lambda: content_generator.generate_from_prompt(
                                user_prompt=slide.prompt,
                                context=context_text,
                                audience=str(project.audience_type),
                                visual_type=slide.visual_type.value,  # Используем .value для получения строки из enum
                                max_chars=200,  # Ограничение для титульного слайда
                            )
                        ),
                        timeout=timeout_sec
                    )
            else:
                out = await asyncio.wait_for(
                    asyncio.to_thread(
                        lambda: content_generator.generate_from_prompt(
                            user_prompt=slide.prompt,
                            context=context_text,
                            audience=str(project.audience_type),
                            visual_type=slide.visual_type.value,  # Используем .value для получения строки из enum
                            max_chars=200,
                        )
                    ),
                    timeout=timeout_sec
                )
            generated_text = (out.get("content") or "").strip()
        except asyncio.TimeoutError:
            error_msg = f"Таймаут генерации титульного слайда ({timeout_sec}с)"
            logger.warning(f"⏱️ {error_msg}")
            generated_text = slide.prompt  # Используем промпт как заголовок
            out = {"title": slide.prompt, "subtitle": "", "layout": "title"}
            job.error_message = error_msg
            await db.commit()
        except Exception as e:
            logger.error(f"Ошибка генерации титульного слайда: {e}")
            generated_text = slide.prompt  # Используем промпт как заголовок
            out = {"title": slide.prompt, "subtitle": "", "layout": "title"}
    else:
        # Генерация может быть очень долгой на CPU/GPU.
        # Чтобы job не "висел" бесконечно в UI, запускаем генерацию в отдельном потоке
        # и ограничиваем её по времени.
        job.progress = 40
        await db.commit()

        try:
            logger.info(f"Начинаем генерацию слайда (таймаут: {timeout_sec}с)")
            # Используем семафор для ограничения параллельных вызовов модели
            # По умолчанию 2 параллельных генерации одновременно
            if _llm_generation_semaphore:
                async with _llm_generation_semaphore:
                    out = await asyncio.wait_for(
                        asyncio.to_thread(
                            lambda: content_generator.generate_from_prompt(
                                user_prompt=slide.prompt,
                                context=context_text,
                                audience=str(project.audience_type),
                                visual_type=slide.visual_type.value,  # Используем .value для получения строки из enum
                                max_chars=800,  # Ограничение символов для слайда
                            )
                        ),
                        timeout=timeout_sec
                    )
            else:
                # Если семафор не инициализирован, работаем без ограничений
                out = await asyncio.wait_for(
                    asyncio.to_thread(
                        lambda: content_generator.generate_from_prompt(
                            user_prompt=slide.prompt,
                            context=context_text,
                            audience=str(project.audience_type),
                            visual_type=slide.visual_type.value,  # Используем .value для получения строки из enum
                            max_chars=800,  # Ограничение символов для слайда
                        )
                    ),
                    timeout=timeout_sec
                )
            generated_text = (out.get("content") or "").strip()
            layout_type = out.get("layout", "title_and_content")
            
            # Логируем сгенерированный контент для отладки
            logger.info(f"Сгенерированный контент для слайда {slide.id} (visual_type={slide.visual_type.value}): {generated_text[:200]}...")
            
            # Генерируем визуализации для таблиц и графиков через matplotlib
            # (переменная generated_image_path уже инициализирована выше)
            if slide.visual_type.value == "table" and generated_text:
                job.progress = 55
                await db.commit()
                try:
                    logger.info(f"Генерация изображения таблицы для слайда {slide.id}")
                    table_image_path = generate_table_image(generated_text, IMAGES_DIR)
                    if table_image_path and Path(table_image_path).exists():
                        generated_image_path = table_image_path
                        logger.info(f"Изображение таблицы успешно сгенерировано: {generated_image_path}")
                    else:
                        logger.warning(f"Изображение таблицы не было создано или файл не найден")
                except Exception as e:
                    logger.warning(f"Не удалось сгенерировать изображение таблицы: {e}", exc_info=True)
            
            elif slide.visual_type.value == "chart" and generated_text:
                job.progress = 55
                await db.commit()
                try:
                    logger.info(f"Генерация изображения графика для слайда {slide.id}")
                    chart_image_path = generate_chart_image(generated_text, IMAGES_DIR)
                    if chart_image_path and Path(chart_image_path).exists():
                        generated_image_path = chart_image_path
                        logger.info(f"Изображение графика успешно сгенерировано: {generated_image_path}")
                    else:
                        logger.warning(f"Изображение графика не было создано или файл не найден")
                except Exception as e:
                    logger.warning(f"Не удалось сгенерировать изображение графика: {e}", exc_info=True)
            
        except asyncio.TimeoutError:
            # При таймауте сохраняем fallback контент вместо полного провала
            error_msg = (
                f"Таймаут генерации ({timeout_sec}с). "
                "Генерация слайда превысила допустимое время. "
                "Попробуйте уменьшить MAX_NEW_TOKENS в настройках или увеличить LLM_GENERATION_TIMEOUT_SEC."
            )
            logger.warning(f"⏱️ {error_msg}")
            # Сохраняем fallback контент вместо полного провала
            generated_text = (
                f"• Таймаут генерации ({timeout_sec}с). "
                "Генерация слайда превысила допустимое время.\n"
                "• Попробуйте:\n"
                "  - Уменьшить MAX_NEW_TOKENS в настройках\n"
                "  - Увеличить LLM_GENERATION_TIMEOUT_SEC\n"
                "  - Использовать более быструю модель или GPU"
            )
            layout_type = "title_and_content"
            # Сохраняем информацию о таймауте в метаданные
            job.error_message = error_msg
            await db.commit()
        
        # Генерируем изображение, если тип визуализации - image
        # Используем описание от Qwen как промпт для Stable Diffusion
        if slide.visual_type.value == "image":
            job.progress = 60
            await db.commit()
            try:
                logger.info(f"Генерация изображения для слайда {slide.id}")
                # Используем сгенерированное описание от Qwen, иначе fallback на промпт пользователя
                image_prompt = (generated_text or "").strip() or slide.prompt.strip()
                prompt_source = "llm" if (generated_text or "").strip() else "user"
                logger.info(
                    f"Промпт для изображения (source={prompt_source}): {image_prompt[:200]}..."
                )
                
                generated_image_path = await image_generator.generate_image_async(
                    prompt=image_prompt,
                )
                if generated_image_path:
                    logger.info(f"Изображение успешно сгенерировано: {generated_image_path}")
                    # Проверяем, что файл существует
                    if not Path(generated_image_path).exists():
                        logger.warning(f"Файл изображения не найден: {generated_image_path}")
                        generated_image_path = None
            except Exception as e:
                logger.warning(f"Не удалось сгенерировать изображение: {e}", exc_info=True)
                # Продолжаем без изображения
                generated_image_path = None

    job.progress = 70
    await db.commit()

    # Сохраняем текстовое содержимое (используется для превью/экспорта)
    sc = SlideContent(slide_id=slide.id)
    _sc_set_text(sc, generated_text)
    
    # Сохраняем метаданные о макете, типе визуализации и пути к изображению
    sc.llm_meta = {
        "layout": layout_type,
        "visual_type": slide.visual_type.value,  # Используем .value для получения строки из enum
    }
    if generated_image_path:
        sc.llm_meta["generated_image_path"] = generated_image_path
    
    # Для титульного слайда сохраняем заголовок и подзаголовок отдельно
    if layout_type == "title" and isinstance(out, dict):
        if "title" in out:
            sc.llm_meta["title"] = out.get("title", "")
        if "subtitle" in out:
            sc.llm_meta["subtitle"] = out.get("subtitle", "")
    # Сохраняем информацию о таймауте, если он произошел
    if job.error_message and "Таймаут" in job.error_message:
        sc.llm_meta["timeout_occurred"] = True
        sc.llm_meta["timeout_seconds"] = timeout_sec

    # если у тебя есть version — увеличим
    if hasattr(sc, "version"):
        res = await db.execute(
            select(SlideContent)
            .where(SlideContent.slide_id == slide.id)
            .order_by(desc(SlideContent.version))
            .limit(1)
        )
        last = res.scalar_one_or_none()
        sc.version = (last.version + 1) if last else 1

    db.add(sc)
    await db.commit()
    
    # Статус слайда будет обновлен в _set_job_done() после завершения job


async def export_project_pptx_job(db, job: Job):
    project = (await db.execute(select(Project).where(Project.id == job.project_id))).scalar_one()

    slides = (await db.execute(
        select(Slide)
        .where(Slide.project_id == project.id)
        .order_by(Slide.position.asc())
    )).scalars().all()

    # шаблон опционален
    template_path = await _resolve_project_template_path(db, project.id)

    # если шаблон назначен, но файла нет — лучше явно упасть (чтобы не молча делать "без шаблона")
    if template_path and not Path(template_path).exists():
        raise FileNotFoundError(f"Template file not found: {template_path}")

    builder = PresentationBuilder(template_path=template_path)

    export_slides: list[SlideExport] = []
    for i, s in enumerate(slides, start=1):
        sc = await _get_latest_slide_content(db, s.id)
        content_text = _sc_get_text(sc)
        
        # Получаем метаданные о макете из llm_meta
        layout_type = "title_and_content"
        images_list = []
        if sc and sc.llm_meta:
            if "layout" in sc.llm_meta:
                layout_type = sc.llm_meta["layout"]
            # Получаем путь к сгенерированному изображению, если есть
            if "generated_image_path" in sc.llm_meta:
                image_path = sc.llm_meta["generated_image_path"]
                if Path(image_path).exists():
                    images_list.append(image_path)
        
        export_slides.append(SlideExport(
            title=s.title or f"Slide {i}", 
            content=content_text, 
            images=images_list,
            layout=layout_type,
            visual_type=s.visual_type.value
        ))

    for se in export_slides:
        builder.add_slide(
            slide_type=se.layout if hasattr(se, "layout") else "content",
            title=se.title, 
            content=se.content, 
            images=se.images,
            visual_type=se.visual_type if hasattr(se, "visual_type") else "text"
        )

    data = builder.save_to_bytes().getvalue()

    safe_title = (project.title or "presentation").strip().replace(" ", "_")
    filename = f"{safe_title}_{uuid.uuid4().hex}.pptx"
    out_path = STORAGE_DIR / filename
    out_path.write_bytes(data)

    out_file = File(
        user_id=job.user_id,
        project_id=project.id,
        kind=FileKind.export_pptx,
        filename=filename,
        storage_path=str(out_path),
        size_bytes=len(data),
    )
    db.add(out_file)
    await db.commit()
    await db.refresh(out_file)

    job.result_file_id = out_file.id
    await db.commit()

async def export_project_pdf_job(db, job: Job):
    project = (await db.execute(select(Project).where(Project.id == job.project_id))).scalar_one()

    slides = (await db.execute(
        select(Slide)
        .where(Slide.project_id == project.id)
        .order_by(Slide.position.asc())
    )).scalars().all()

    export_slides: list[SlideExport] = []
    for i, s in enumerate(slides, start=1):
        sc = await _get_latest_slide_content(db, s.id)
        content_text = _sc_get_text(sc)
        
        # Получаем метаданные о макете из llm_meta
        layout_type = "title_and_content"
        images_list = []
        if sc and sc.llm_meta:
            if "layout" in sc.llm_meta:
                layout_type = sc.llm_meta["layout"]
            # Получаем путь к сгенерированному изображению, если есть
            if "generated_image_path" in sc.llm_meta:
                image_path = sc.llm_meta["generated_image_path"]
                if Path(image_path).exists():
                    images_list.append(image_path)
        
        export_slides.append(SlideExport(
            title=s.title or f"Slide {i}", 
            content=content_text, 
            images=images_list,
            layout=layout_type,
            visual_type=s.visual_type.value
        ))

    pdf_bytes = slides_to_pdf_bytes(export_slides, audience=str(project.audience_type))

    safe_title = (project.title or "presentation").strip().replace(" ", "_")
    filename = f"{safe_title}_{uuid.uuid4().hex}.pdf"
    out_path = STORAGE_DIR / filename
    out_path.write_bytes(pdf_bytes)

    out_file = File(
        user_id=job.user_id,
        project_id=project.id,
        kind=FileKind.export_pdf,
        filename=filename,
        storage_path=str(out_path),
        size_bytes=len(pdf_bytes),
    )
    db.add(out_file)
    await db.commit()
    await db.refresh(out_file)

    job.result_file_id = out_file.id
    await db.commit()


async def handle_job(db: AsyncSession, job: Job) -> None:
    if job.type == JobType.generate_slide:
        await generate_slide_job(db, job)
        return

    if job.type == JobType.export_pptx:
        await export_project_pptx_job(db, job)
        return

    if job.type == JobType.export_pdf:
        await export_project_pdf_job(db, job)
        return

    raise NotImplementedError(f"Job type not supported yet: {job.type}")


# ----------------------------
# Job state helpers
# ----------------------------
async def _fetch_one_queued_job(db: AsyncSession) -> Optional[Job]:
    q = (
        select(Job)
        .where(Job.status == JobStatus.queued)
        .order_by(Job.id.asc())
        .with_for_update(skip_locked=True)
        .limit(1)
    )
    res = await db.execute(q)
    return res.scalar_one_or_none()


async def _fetch_multiple_queued_jobs(db: AsyncSession, limit: int) -> List[Job]:
    """Получает несколько задач из очереди для параллельной обработки"""
    q = (
        select(Job)
        .where(Job.status == JobStatus.queued)
        .order_by(Job.id.asc())
        .with_for_update(skip_locked=True)
        .limit(limit)
    )
    res = await db.execute(q)
    return list(res.scalars().all())


async def _set_job_running(db: AsyncSession, job: Job) -> None:
    job.status = JobStatus.running
    job.progress = 1
    job.error_message = None
    await db.commit()


async def _set_job_done(db: AsyncSession, job: Job) -> None:
    job.status = JobStatus.done
    job.progress = 100
    
    # Обновляем статус слайда на ready только после того, как job помечен как done
    # Это гарантирует, что "Готов" появится только после полного завершения генерации (после лога "Job done")
    if job.type == JobType.generate_slide and job.slide_id is not None:
        slide = (await db.execute(select(Slide).where(Slide.id == job.slide_id))).scalar_one_or_none()
        if slide is not None and hasattr(slide, "status"):
            slide.status = SlideStatus.ready
    
    await db.commit()


async def _set_job_failed(db: AsyncSession, job: Job, exc: Exception) -> None:
    job.status = JobStatus.error
    job.progress = 100
    job.error_message = str(exc)
    # если это job генерации — отметим и слайд как error, чтобы UI перестал крутиться
    try:
        if job.slide_id is not None:
            slide = (await db.execute(select(Slide).where(Slide.id == job.slide_id))).scalar_one_or_none()
            if slide is not None and hasattr(slide, "status"):
                slide.status = SlideStatus.error
    except Exception:
        # не мешаем падению job
        pass
    await db.commit()


# ----------------------------
# Worker task handler
# ----------------------------
async def _process_job(job_id: int) -> None:
    """Обрабатывает одну задачу в отдельной сессии БД"""
    async with AsyncSessionLocal() as db:
        try:
            # Загружаем задачу из БД в текущей сессии
            job = (await db.execute(select(Job).where(Job.id == job_id))).scalar_one_or_none()
            if not job:
                logger.warning(f"Job {job_id} not found")
                return
            
            # Обновляем статус задачи
            await _set_job_running(db, job)
            
            logger.info(f"➡️ Processing job id={job.id} type={job.type}")
            
            # Обрабатываем задачу (семафор применяется внутри handle_job для генерации слайдов)
            await handle_job(db, job)
            
            # Отмечаем как выполненную
            await _set_job_done(db, job)
            logger.info(f"✅ Job done id={job.id}")
        except Exception as e:
            logger.exception(f"❌ Job failed id={job_id}: {e}")
            try:
                # Перезагружаем задачу для обновления статуса
                job = (await db.execute(select(Job).where(Job.id == job_id))).scalar_one_or_none()
                if job:
                    await _set_job_failed(db, job, e)
            except Exception as db_error:
                logger.error(f"Failed to update job status in DB: {db_error}")


# ----------------------------
# Worker loop
# ----------------------------
async def worker_loop() -> None:
    global _llm_generation_semaphore
    
    parallel_jobs = settings.WORKER_PARALLEL_JOBS
    # Инициализируем семафор для ограничения параллельных вызовов модели
    # По умолчанию 2 параллельных генерации одновременно
    llm_parallel_limit = int(os.getenv("LLM_PARALLEL_GENERATIONS", "2"))
    _llm_generation_semaphore = asyncio.Semaphore(llm_parallel_limit)
    
    logger.info(f"🧵 Worker started (parallel jobs: {parallel_jobs}, LLM parallel limit: {llm_parallel_limit})")
    
    active_tasks: set[asyncio.Task] = set()

    while True:
        # Убираем завершенные задачи
        active_tasks = {t for t in active_tasks if not t.done()}
        
        # Сколько задач можем запустить параллельно
        available_slots = parallel_jobs - len(active_tasks)
        
        if available_slots > 0:
            # Получаем задачи из очереди
            async with AsyncSessionLocal() as db:
                jobs = await _fetch_multiple_queued_jobs(db, limit=available_slots)
            
            # Запускаем обработку задач параллельно
            for job in jobs:
                task = asyncio.create_task(_process_job(job.id))
                active_tasks.add(task)
                logger.info(f"🚀 Started parallel processing job id={job.id}")
        
        # Ждем немного перед следующей проверкой
        if not active_tasks:
            await asyncio.sleep(POLL_INTERVAL_SEC)
        else:
            # Ждем завершения хотя бы одной задачи
            done, pending = await asyncio.wait(active_tasks, return_when=asyncio.FIRST_COMPLETED, timeout=POLL_INTERVAL_SEC)
            active_tasks = pending


def main() -> None:
    asyncio.run(worker_loop())


if __name__ == "__main__":
    main()

