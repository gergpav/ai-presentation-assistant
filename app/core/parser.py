from pathlib import Path
from starlette.datastructures import UploadFile as StarletteUploadFile
from fastapi import UploadFile
import io
import pandas as pd
from docx import Document
import pdfplumber
from pptx import Presentation   # <--- добавили импорт для PPTX
from typing import Dict, Any
import logging

logger = logging.getLogger(__name__)


async def extract_text_from_file(file: UploadFile) -> Dict[str, Any]:
    """Извлекает текст и структурированные данные из файла"""
    content = await file.read()
    filename = file.filename.lower()

    result = {
        "text": "",
        "tables": [],
        "metadata": {"filename": file.filename, "type": filename.split('.')[-1]}
    }

    try:
        # ---- TXT УБРАН, .txt больше не поддерживаем ----
        # if filename.endswith(".txt"):
        #     result["text"] = content.decode("utf-8", errors="ignore")

        # ---- DOCX ----
        if filename.endswith(".docx"):
            doc = Document(io.BytesIO(content))
            result["text"] = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

            # Извлекаем таблицы из DOCX
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    result["tables"].append(table_data)

        # ---- PDF ----
        elif filename.endswith(".pdf"):
            text = ""
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text() or ""
                    text += page_text + "\n"

                    # Извлекаем таблицы из PDF
                    tables = page.extract_tables()
                    for table in tables:
                        if table and any(any(cell is not None for cell in row) for row in table):
                            result["tables"].append(table)

            result["text"] = text.strip()

        # ---- XLSX (исправленная работа с несколькими листами) ----
        elif filename.endswith(".xlsx"):
            # Обработка Excel файлов
            excel_file = io.BytesIO(content)

            # Читаем все листы через ExcelFile, чтобы не гонять указатель BytesIO
            xl = pd.ExcelFile(excel_file)
            all_text = f"Excel файл: {file.filename}\n"
            all_tables = []

            for sheet_name in xl.sheet_names:
                # ключевое изменение: читаем через xl.parse, а не pd.read_excel(excel_file, ...)
                df = xl.parse(sheet_name=sheet_name)

                all_text += f"\n--- Лист: {sheet_name} ---\n"
                all_text += df.to_string() + "\n"

                # Сохраняем таблицу как список списков
                all_tables.append({
                    "sheet_name": sheet_name,
                    "data": df.fillna("").values.tolist(),
                    "columns": df.columns.tolist()
                })

            result["text"] = all_text
            result["tables"] = all_tables

        # ---- PPTX (новая обработка презентаций) ----
        elif filename.endswith(".pptx"):
            prs = Presentation(io.BytesIO(content))
            all_text_parts = []
            all_tables = []

            for slide_idx, slide in enumerate(prs.slides):
                # Заголовок слайда
                if slide.shapes.title and slide.shapes.title.text:
                    title = slide.shapes.title.text.strip()
                    if title:
                        all_text_parts.append(f"Слайд {slide_idx + 1} — заголовок: {title}")

                # Остальной текст и таблицы
                for shape in slide.shapes:
                    # Текстовые элементы
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            all_text_parts.append(f"Слайд {slide_idx + 1}: {text}")

                    # Таблицы
                    if hasattr(shape, "has_table") and shape.has_table:
                        table = shape.table
                        table_data = []
                        for row in table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append(row_data)
                        if table_data:
                            all_tables.append({
                                "slide_index": slide_idx,
                                "data": table_data
                            })

            result["text"] = "\n".join(all_text_parts)
            result["tables"] = all_tables

        else:
            # Всё остальное считаем неподдерживаемым форматом
            raise ValueError(f"Неподдерживаемый формат файла: {filename}")

        logger.info(
            f"Успешно обработан файл {file.filename}: "
            f"{len(result['text'])} символов, {len(result['tables'])} таблиц"
        )
        return result

    except Exception as e:
        logger.error(f"Ошибка обработки файла {file.filename}: {e}")
        raise ValueError(f"Ошибка обработки файла: {str(e)}")


async def extract_text_from_path(path: str) -> Dict[str, Any]:
    """
    Мостик для воркера: читаем файл с диска и прогоняем через существующий extract_text_from_file.
    """
    p = Path(path)
    content = p.read_bytes()

    # UploadFile ожидает file-like объект
    file_like = io.BytesIO(content)
    uf = StarletteUploadFile(filename=p.name, file=file_like)
    return await extract_text_from_file(uf)

