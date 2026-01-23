# app/core/pptx_builder.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE
from pptx.dml.color import RGBColor
import io
import logging
from typing import List, Optional
import re

logger = logging.getLogger(__name__)


class PresentationBuilder:
    def __init__(self, template_path: Optional[str] = None):
        if template_path:
            self.prs = Presentation(template_path)
            self.layout_map = {}
            for i, layout in enumerate(self.prs.slide_layouts):
                name = layout.name.lower()
                if 'title and content' in name:
                    self.layout_map['title_and_content'] = i
                elif 'two content' in name:
                    self.layout_map['two_content'] = i
                elif 'title only' in name:
                    self.layout_map['title_only'] = i
                elif 'title slide' in name:
                    self.layout_map['title'] = i
            self.clear_slides()
        else:
            self.prs = Presentation()

    def clear_slides(self):
        """Очищает все слайды из шаблона, но сохраняет дизайн/layout'ы"""
        try:
            sldIdLst = self.prs.slides._sldIdLst
            for sldId in list(sldIdLst):
                sldIdLst.remove(sldId)
            logger.info("✅ Шаблон очищен")
        except Exception as e:
            logger.warning(f"Не удалось очистить шаблон: {e}")
            self.prs = Presentation()

    def _get_layout_index(self, layout_type: str) -> int:
        """Получает индекс макета по его типу"""
        if layout_type in self.layout_map:
            return self.layout_map[layout_type]
        
        # Fallback: пытаемся найти по имени
        for i, layout in enumerate(self.prs.slide_layouts):
            name = layout.name.lower()
            if layout_type == "title" and 'title slide' in name:
                return i
            elif layout_type == "two_content" and 'two content' in name:
                return i
            elif layout_type == "title_only" and 'title only' in name:
                return i
            elif layout_type == "title_and_content" and 'title and content' in name:
                return i
        
        # Если ничего не найдено, используем первый доступный макет
        return 0

    def add_slide(self, slide_type: str, title: str, content: str,
                  images: Optional[List[str]] = None, visual_type: str = "text"):
        """
        Add a new slide to the presentation with support for different layouts and visualization types.

        Parameters:
        - slide_type: "title", "two_content", "title_only", "title_and_content"
        - title: заголовок слайда
        - content: содержимое слайда (текст, таблица, данные для графика)
        - images: список путей к изображениям
        - visual_type: "text", "table", "chart", "image"
        """
        images = images or []
        
        # Определяем индекс макета
        layout_idx = self._get_layout_index(slide_type)

        try:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
        except Exception as e:
            logger.warning(f"Не удалось использовать макет {slide_type}, используем первый доступный: {e}")
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])

        # Заголовок
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title
            # Применяем стили к заголовку
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(44)  # Стандартный размер заголовка
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.LEFT

        # Обработка контента в зависимости от типа визуализации
        if slide_type == "title" or slide_type == "title_only" or not content.strip():
            # Для титульного слайда и слайда только с заголовком ничего не добавляем
            pass
        elif slide_type == "two_content":
            # Для двухколоночного слайда разделяем контент на две части
            self._add_two_content(slide, content, visual_type)
        elif visual_type == "table":
            self._add_table(slide, content)
        elif visual_type == "chart":
            self._add_chart(slide, content)
        elif visual_type == "image":
            # Для изображений: если есть изображения, добавляем их, иначе добавляем описание
            if images:
                self._add_images(slide, images)
            elif content.strip():
                # Если изображений нет, добавляем описание как текст
                self._add_text_content(slide, content)
        else:  # text
            self._add_text_content(slide, content)

        # Картинки (если есть и тип не image - добавляем дополнительно)
        if images and visual_type != "image":
            self._add_images(slide, images)

        return slide
    
    def _add_text_content(self, slide, content: str):
        """Добавляет текстовый контент на слайд с правильным форматированием"""
        body_shape = None
        try:
            for ph in slide.placeholders:
                if ph.placeholder_format.type == PP_PLACEHOLDER.BODY:
                    body_shape = ph
                    break
        except Exception:
            body_shape = None

        if body_shape is not None:
            tf = body_shape.text_frame
            tf.clear()
            tf.word_wrap = True
            
            # Разбиваем контент на строки и создаем параграфы
            lines = content.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # Убираем маркеры списка если есть
                clean_line = line.strip()
                if clean_line.startswith('•'):
                    clean_line = clean_line[1:].strip()
                elif clean_line.startswith('-'):
                    clean_line = clean_line[1:].strip()
                
                p.text = clean_line
                p.font.size = Pt(18)
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(6)
                p.level = 0
                
                # Определяем жирный текст по заглавным буквам или ключевым словам
                if clean_line.isupper() or any(keyword in clean_line.lower() for keyword in ['важно', 'ключевой', 'основной']):
                    p.font.bold = True
        else:
            # Фоллбек: текстбокс
            left = Inches(1)
            top = Inches(2)
            width = Inches(11)
            height = Inches(4)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True
            tf.text = content
            for p in tf.paragraphs:
                p.font.size = Pt(18)
                p.alignment = PP_ALIGN.LEFT
    
    def _add_table(self, slide, content: str):
        """Добавляет таблицу на слайд из текстового контента"""
        # Парсим таблицу из текста (формат: Столбец1 | Столбец2 | ...)
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        if not lines:
            return
        
        # Улучшенный парсинг: ищем строки с разделителями |
        table_rows = []
        for line in lines:
            # Убираем маркеры списка и лишние символы
            clean_line = line.strip()
            if clean_line.startswith('•') or clean_line.startswith('-'):
                clean_line = clean_line[1:].strip()
            
            # Проверяем, есть ли разделитель |
            if '|' in clean_line:
                cells = [col.strip() for col in clean_line.split('|')]
                if len(cells) >= 2:  # Минимум 2 столбца
                    table_rows.append(cells)
        
        if not table_rows or len(table_rows) < 2:
            # Если не удалось распарсить таблицу, добавляем как текст
            logger.warning("Не удалось распарсить таблицу, добавляем как текст")
            self._add_text_content(slide, content)
            return
        
        # Определяем количество столбцов по максимальному количеству в строках
        num_cols = max(len(row) for row in table_rows)
        num_rows = len(table_rows)
        
        # Создаем таблицу
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(11)
        height = Inches(4.5)
        
        try:
            table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
            table = table_shape.table
            
            # Заполняем таблицу
            for row_idx, row_data in enumerate(table_rows):
                for col_idx in range(num_cols):
                    cell = table.cell(row_idx, col_idx)
                    if col_idx < len(row_data):
                        cell.text = row_data[col_idx]
                    else:
                        cell.text = ""
                    
                    # Стили для заголовка (первая строка)
                    if row_idx == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)  # Синий
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = RGBColor(255, 255, 255)  # Белый
                            paragraph.font.size = Pt(14)
                            paragraph.alignment = PP_ALIGN.CENTER
                    else:
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(12)
                            paragraph.alignment = PP_ALIGN.LEFT
        except Exception as e:
            logger.warning(f"Не удалось создать таблицу: {e}, добавляем как текст")
            self._add_text_content(slide, content)
    
    def _add_chart(self, slide, content: str):
        """Добавляет реальный график на слайд из данных"""
        # Парсим данные графика (формат: Название: Значение)
        chart_data = {}
        lines = [line.strip() for line in content.split('\n') if line.strip() and ':' in line]
        
        for line in lines:
            # Убираем маркеры списка
            clean_line = line.strip()
            if clean_line.startswith('•') or clean_line.startswith('-'):
                clean_line = clean_line[1:].strip()
            
            parts = clean_line.split(':', 1)
            if len(parts) == 2:
                name = parts[0].strip()
                value_str = parts[1].strip()
                # Пытаемся извлечь числовое значение
                try:
                    # Убираем все нечисловые символы кроме точки и минуса
                    numeric_str = ''.join(c for c in value_str if c.isdigit() or c == '.' or c == '-')
                    if numeric_str:
                        value = float(numeric_str)
                        chart_data[name] = value
                except ValueError:
                    # Если не число, используем как есть
                    chart_data[name] = value_str
        
        if not chart_data:
            # Если не удалось распарсить данные, добавляем как текст
            self._add_text_content(slide, content)
            return
        
        try:
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
            
            # Создаем данные для графика
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = list(chart_data.keys())
            chart_data_obj.add_series('Значения', list(chart_data.values()))
            
            # Определяем позицию и размер графика
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(11)
            height = Inches(5)
            
            # Создаем график (столбчатый)
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                left, top, width, height,
                chart_data_obj
            )
            
            # Настраиваем стиль графика
            chart = chart_shape.chart
            chart.has_legend = True
            chart.has_title = True
            chart.chart_title.text_frame.text = "График данных"
            
        except Exception as e:
            logger.warning(f"Не удалось создать график: {e}, добавляем как текст")
            # Форматируем как список для отображения
            formatted_lines = []
            for name, value in chart_data.items():
                formatted_lines.append(f"• {name}: {value}")
            formatted_content = "\n".join(formatted_lines)
            self._add_text_content(slide, formatted_content)
    
    def _add_two_content(self, slide, content: str, visual_type: str = "text"):
        """Добавляет контент в две колонки для слайда типа two_content"""
        # Разделяем контент на две части (по разделителю или пополам)
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        
        # Ищем разделитель для двух колонок
        separator_keywords = ["---", "===", "раздел", "часть"]
        split_idx = None
        
        for i, line in enumerate(lines):
            if any(sep in line.lower() for sep in separator_keywords):
                split_idx = i
                break
        
        # Если разделитель не найден, делим пополам
        if split_idx is None:
            split_idx = len(lines) // 2
        
        left_content = "\n".join(lines[:split_idx])
        right_content = "\n".join(lines[split_idx:])
        
        # Ищем placeholders для двух колонок
        left_shape = None
        right_shape = None
        
        try:
            for ph in slide.placeholders:
                ph_type = ph.placeholder_format.type
                if ph_type == PP_PLACEHOLDER.BODY:
                    if left_shape is None:
                        left_shape = ph
                    elif right_shape is None:
                        right_shape = ph
                        break
        except Exception:
            pass
        
        # Если нашли два placeholder'а, используем их
        if left_shape and right_shape:
            # Левая колонка
            tf_left = left_shape.text_frame
            tf_left.clear()
            tf_left.word_wrap = True
            for i, line in enumerate(left_content.split('\n')):
                if i == 0:
                    p = tf_left.paragraphs[0]
                else:
                    p = tf_left.add_paragraph()
                clean_line = line.strip()
                if clean_line.startswith('•') or clean_line.startswith('-'):
                    clean_line = clean_line[1:].strip()
                p.text = clean_line
                p.font.size = Pt(16)
                p.alignment = PP_ALIGN.LEFT
            
            # Правая колонка
            tf_right = right_shape.text_frame
            tf_right.clear()
            tf_right.word_wrap = True
            for i, line in enumerate(right_content.split('\n')):
                if i == 0:
                    p = tf_right.paragraphs[0]
                else:
                    p = tf_right.add_paragraph()
                clean_line = line.strip()
                if clean_line.startswith('•') or clean_line.startswith('-'):
                    clean_line = clean_line[1:].strip()
                p.text = clean_line
                p.font.size = Pt(16)
                p.alignment = PP_ALIGN.LEFT
        else:
            # Fallback: создаем два текстбокса вручную
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(5.5)
            height = Inches(5)
            
            # Левая колонка
            textbox_left = slide.shapes.add_textbox(left, top, width, height)
            tf_left = textbox_left.text_frame
            tf_left.word_wrap = True
            tf_left.text = left_content
            for p in tf_left.paragraphs:
                p.font.size = Pt(16)
                p.alignment = PP_ALIGN.LEFT
            
            # Правая колонка
            textbox_right = slide.shapes.add_textbox(left + width + Inches(0.3), top, width, height)
            tf_right = textbox_right.text_frame
            tf_right.word_wrap = True
            tf_right.text = right_content
            for p in tf_right.paragraphs:
                p.font.size = Pt(16)
                p.alignment = PP_ALIGN.LEFT
    
    def _add_images(self, slide, images: List[str]):
        """Добавляет изображения на слайд"""
        if not images:
            return
        
        # Для одного изображения - размещаем по центру, большего размера
        if len(images) == 1:
            img_path = images[0]
            try:
                # Проверяем существование файла
                if not Path(img_path).exists():
                    logger.warning(f"Изображение не найдено: {img_path}")
                    return
                
                # Размещаем по центру слайда
                left = Inches(1.5)
                top = Inches(2)
                width = Inches(9)
                height = Inches(5)
                slide.shapes.add_picture(img_path, left, top, width=width, height=height)
            except Exception as e:
                logger.warning(f"Не удалось добавить изображение '{img_path}': {e}")
        else:
            # Для нескольких изображений - размещаем в сетке
            left0 = Inches(1)
            top0 = Inches(6.0)
            cell_w = Inches(5.5)
            cell_h = Inches(1.6)
            max_imgs = min(4, len(images))

            for i, img_path in enumerate(images[:max_imgs]):
                if not Path(img_path).exists():
                    logger.warning(f"Изображение не найдено: {img_path}")
                    continue
                
                row, col = divmod(i, 2)
                left = left0 + col * (cell_w + Inches(0.3))
                top = top0 + row * (cell_h + Inches(0.3))
                try:
                    slide.shapes.add_picture(img_path, left, top, width=cell_w, height=cell_h)
                except Exception as e:
                    logger.warning(f"Не удалось добавить изображение '{img_path}': {e}")

    def save_to_bytes(self) -> io.BytesIO:
        bytes_io = io.BytesIO()
        self.prs.save(bytes_io)
        bytes_io.seek(0)
        return bytes_io

    def get_slide_count(self) -> int:
        return len(self.prs.slides)



