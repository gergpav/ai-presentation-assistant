from io import BytesIO

from pptx import Presentation
from typing import List, Dict
from pydantic import BaseModel
import logging

logger = logging.getLogger(__name__)


class TemplateAnalysis(BaseModel):
    slides_count: int
    layouts: List[Dict]


def analyze_template(template_bytes: str) -> TemplateAnalysis:
    """Анализирует PowerPoint шаблон и возвращает информацию о макетах"""
    try:
        prs = Presentation(BytesIO(template_bytes))
        layouts = []

        for i, layout in enumerate(prs.slide_layouts):
            layouts.append({
                "index": i,
                "name": layout.name,
                "placeholders": [
                    {
                        "type": "text",
                        "index": ph.placeholder_format.idx
                    }
                    for ph in layout.placeholders
                    if ph.has_text_frame
                ]
            })

        return TemplateAnalysis(
            slides_count=len(prs.slides),
            layouts=layouts
        )
    except Exception as e:
        logger.error(f"Ошибка анализа шаблона: {e}")
        raise